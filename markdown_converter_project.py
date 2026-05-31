# markdown_converter 项目模板
# 说明：下面是一个 Codex-ready 的多文件项目模板。
# 使用时按文件路径拆分保存即可。

# ============================================================
# file: pyproject.toml
# ============================================================
"""
[project]
name = "markdown-converter"
version = "0.1.0"
description = "Convert office documents and files to Markdown"
requires-python = ">=3.9"
dependencies = [
    "python-docx>=1.1.0",
    "openpyxl>=3.1.0",
    "pymupdf>=1.24.0",
    "pdfplumber>=0.11.0",
    "python-pptx>=0.6.23",
    "click>=8.1.0",
]

[project.scripts]
mdconvert = "markdown_converter.cli:main"

[tool.setuptools.packages.find]
where = ["src"]
"""

# ============================================================
# file: README.md
# ============================================================
"""
# markdown-converter

一个用于将常见文件和办公文档转换为 Markdown 的 Python 工具。

## 支持格式

- `.txt`
- `.md`
- `.docx`
- `.xlsx`
- `.pdf`
- `.pptx`

## 安装

```bash
pip install -e .
```

## 使用

转换单个文件：

```bash
mdconvert input.docx -o output
```

批量转换目录：

```bash
mdconvert ./docs -o ./markdown --recursive
```

导出图片资源：

```bash
mdconvert report.docx -o output --assets
```

## 设计原则

- 每种文件类型一个 converter
- 统一返回 Markdown 字符串
- 表格转换为 Markdown table
- 图片导出到 assets 目录，Markdown 中使用相对路径
- PDF 优先保留文本，其次尝试提取表格
"""

# ============================================================
# file: src/markdown_converter/__init__.py
# ============================================================
__version__ = "0.1.0"

# ============================================================
# file: src/markdown_converter/config.py
# ============================================================
from dataclasses import dataclass
from pathlib import Path


@dataclass
class ConvertOptions:
    output_dir: Path
    export_assets: bool = False
    recursive: bool = False
    overwrite: bool = True

# ============================================================
# file: src/markdown_converter/utils.py
# ============================================================
from __future__ import annotations

import re
from pathlib import Path
from typing import Iterable, List, Sequence


def normalize_text(text: str | None) -> str:
    if not text:
        return ""
    text = text.replace("\r\n", "\n").replace("\r", "\n")
    text = re.sub(r"\n{3,}", "\n\n", text)
    return text.strip()


def escape_md_cell(value: object) -> str:
    text = "" if value is None else str(value)
    text = text.replace("\n", "<br>").replace("|", "\\|")
    return text.strip()


def table_to_markdown(rows: Sequence[Sequence[object]]) -> str:
    cleaned_rows: List[List[str]] = [
        [escape_md_cell(cell) for cell in row]
        for row in rows
        if any(str(cell or "").strip() for cell in row)
    ]

    if not cleaned_rows:
        return ""

    max_cols = max(len(row) for row in cleaned_rows)
    normalized = [row + [""] * (max_cols - len(row)) for row in cleaned_rows]

    header = normalized[0]
    separator = ["---"] * max_cols
    body = normalized[1:]

    lines = [
        "| " + " | ".join(header) + " |",
        "| " + " | ".join(separator) + " |",
    ]

    for row in body:
        lines.append("| " + " | ".join(row) + " |")

    return "\n".join(lines)


def ensure_output_dir(path: Path) -> None:
    path.mkdir(parents=True, exist_ok=True)


def write_markdown(content: str, output_path: Path, overwrite: bool = True) -> Path:
    if output_path.exists() and not overwrite:
        raise FileExistsError(f"Output file already exists: {output_path}")
    output_path.parent.mkdir(parents=True, exist_ok=True)
    output_path.write_text(content.strip() + "\n", encoding="utf-8")
    return output_path


def iter_supported_files(input_path: Path, extensions: set[str], recursive: bool = False) -> Iterable[Path]:
    if input_path.is_file():
        if input_path.suffix.lower() in extensions:
            yield input_path
        return

    pattern = "**/*" if recursive else "*"
    for path in input_path.glob(pattern):
        if path.is_file() and path.suffix.lower() in extensions:
            yield path


def safe_asset_name(source: Path, suffix: str, index: int) -> str:
    stem = re.sub(r"[^a-zA-Z0-9_-]+", "_", source.stem).strip("_") or "asset"
    return f"{stem}_{index}{suffix}"

# ============================================================
# file: src/markdown_converter/converters/base.py
# ============================================================
from abc import ABC, abstractmethod
from pathlib import Path
from markdown_converter.config import ConvertOptions


class BaseConverter(ABC):
    @abstractmethod
    def convert(self, input_path: Path, options: ConvertOptions) -> str:
        raise NotImplementedError

# ============================================================
# file: src/markdown_converter/converters/text_converter.py
# ============================================================
from pathlib import Path
from markdown_converter.config import ConvertOptions
from markdown_converter.converters.base import BaseConverter
from markdown_converter.utils import normalize_text


class TextConverter(BaseConverter):
    def convert(self, input_path: Path, options: ConvertOptions) -> str:
        return normalize_text(input_path.read_text(encoding="utf-8", errors="replace"))

# ============================================================
# file: src/markdown_converter/converters/docx_converter.py
# ============================================================
from pathlib import Path
from docx import Document
from markdown_converter.config import ConvertOptions
from markdown_converter.converters.base import BaseConverter
from markdown_converter.utils import normalize_text, table_to_markdown


class DocxConverter(BaseConverter):
    def convert(self, input_path: Path, options: ConvertOptions) -> str:
        doc = Document(str(input_path))
        parts: list[str] = []

        for block in self._iter_blocks(doc):
            kind, value = block
            if kind == "paragraph":
                md = self._paragraph_to_markdown(value)
                if md:
                    parts.append(md)
            elif kind == "table":
                table_md = table_to_markdown(value)
                if table_md:
                    parts.append(table_md)

        return normalize_text("\n\n".join(parts))

    def _iter_blocks(self, doc: Document):
        # 简化版本：先输出所有段落，再输出所有表格。
        # 如需完全保持原文档顺序，可进一步解析 XML body。
        for paragraph in doc.paragraphs:
            yield "paragraph", paragraph
        for table in doc.tables:
            rows = [[cell.text for cell in row.cells] for row in table.rows]
            yield "table", rows

    def _paragraph_to_markdown(self, paragraph) -> str:
        text = paragraph.text.strip()
        if not text:
            return ""

        style_name = (paragraph.style.name or "").lower()

        if "heading 1" in style_name:
            return f"# {text}"
        if "heading 2" in style_name:
            return f"## {text}"
        if "heading 3" in style_name:
            return f"### {text}"
        if "heading 4" in style_name:
            return f"#### {text}"
        if "list bullet" in style_name:
            return f"- {text}"
        if "list number" in style_name:
            return f"1. {text}"

        # 保留 run 级别的粗体/斜体。
        chunks: list[str] = []
        for run in paragraph.runs:
            run_text = run.text
            if not run_text:
                continue
            if run.bold and run.italic:
                run_text = f"***{run_text}***"
            elif run.bold:
                run_text = f"**{run_text}**"
            elif run.italic:
                run_text = f"*{run_text}*"
            chunks.append(run_text)

        return "".join(chunks).strip() or text

# ============================================================
# file: src/markdown_converter/converters/xlsx_converter.py
# ============================================================
from pathlib import Path
from openpyxl import load_workbook
from markdown_converter.config import ConvertOptions
from markdown_converter.converters.base import BaseConverter
from markdown_converter.utils import normalize_text, table_to_markdown


class XlsxConverter(BaseConverter):
    def convert(self, input_path: Path, options: ConvertOptions) -> str:
        wb = load_workbook(filename=str(input_path), data_only=True)
        parts: list[str] = []

        for ws in wb.worksheets:
            parts.append(f"# Sheet: {ws.title}")
            rows = []
            for row in ws.iter_rows(values_only=True):
                if any(cell is not None and str(cell).strip() for cell in row):
                    rows.append(list(row))
            table_md = table_to_markdown(rows)
            if table_md:
                parts.append(table_md)

        return normalize_text("\n\n".join(parts))

# ============================================================
# file: src/markdown_converter/converters/pdf_converter.py
# ============================================================
from pathlib import Path
import fitz  # PyMuPDF
import pdfplumber
from markdown_converter.config import ConvertOptions
from markdown_converter.converters.base import BaseConverter
from markdown_converter.utils import normalize_text, table_to_markdown


class PdfConverter(BaseConverter):
    def convert(self, input_path: Path, options: ConvertOptions) -> str:
        parts: list[str] = []

        text = self._extract_text_with_pymupdf(input_path)
        if text:
            parts.append(text)

        tables = self._extract_tables_with_pdfplumber(input_path)
        if tables:
            parts.append("# Tables")
            parts.extend(tables)

        return normalize_text("\n\n".join(parts))

    def _extract_text_with_pymupdf(self, input_path: Path) -> str:
        doc = fitz.open(str(input_path))
        pages: list[str] = []
        for index, page in enumerate(doc, start=1):
            page_text = page.get_text("text").strip()
            if page_text:
                pages.append(f"## Page {index}\n\n{page_text}")
        return "\n\n".join(pages)

    def _extract_tables_with_pdfplumber(self, input_path: Path) -> list[str]:
        results: list[str] = []
        with pdfplumber.open(str(input_path)) as pdf:
            for page_index, page in enumerate(pdf.pages, start=1):
                for table_index, table in enumerate(page.extract_tables() or [], start=1):
                    table_md = table_to_markdown(table)
                    if table_md:
                        results.append(f"## Page {page_index}, Table {table_index}\n\n{table_md}")
        return results

# ============================================================
# file: src/markdown_converter/converters/pptx_converter.py
# ============================================================
from pathlib import Path
from pptx import Presentation
from markdown_converter.config import ConvertOptions
from markdown_converter.converters.base import BaseConverter
from markdown_converter.utils import normalize_text, table_to_markdown


class PptxConverter(BaseConverter):
    def convert(self, input_path: Path, options: ConvertOptions) -> str:
        prs = Presentation(str(input_path))
        parts: list[str] = []

        for slide_index, slide in enumerate(prs.slides, start=1):
            parts.append(f"# Slide {slide_index}")
            slide_text: list[str] = []

            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_text.append(shape.text.strip())

                if getattr(shape, "has_table", False):
                    rows = []
                    for row in shape.table.rows:
                        rows.append([cell.text for cell in row.cells])
                    table_md = table_to_markdown(rows)
                    if table_md:
                        slide_text.append(table_md)

            if slide_text:
                parts.append("\n\n".join(slide_text))

        return normalize_text("\n\n".join(parts))

# ============================================================
# file: src/markdown_converter/converter_registry.py
# ============================================================
from pathlib import Path
from markdown_converter.converters.base import BaseConverter
from markdown_converter.converters.text_converter import TextConverter
from markdown_converter.converters.docx_converter import DocxConverter
from markdown_converter.converters.xlsx_converter import XlsxConverter
from markdown_converter.converters.pdf_converter import PdfConverter
from markdown_converter.converters.pptx_converter import PptxConverter


CONVERTERS: dict[str, BaseConverter] = {
    ".txt": TextConverter(),
    ".md": TextConverter(),
    ".docx": DocxConverter(),
    ".xlsx": XlsxConverter(),
    ".pdf": PdfConverter(),
    ".pptx": PptxConverter(),
}

SUPPORTED_EXTENSIONS = set(CONVERTERS.keys())


def get_converter(input_path: Path) -> BaseConverter:
    ext = input_path.suffix.lower()
    try:
        return CONVERTERS[ext]
    except KeyError as exc:
        supported = ", ".join(sorted(SUPPORTED_EXTENSIONS))
        raise ValueError(f"Unsupported file type: {ext}. Supported: {supported}") from exc

# ============================================================
# file: src/markdown_converter/service.py
# ============================================================
from pathlib import Path
from markdown_converter.config import ConvertOptions
from markdown_converter.converter_registry import SUPPORTED_EXTENSIONS, get_converter
from markdown_converter.utils import ensure_output_dir, iter_supported_files, write_markdown


class MarkdownConvertService:
    def convert_path(self, input_path: Path, options: ConvertOptions) -> list[Path]:
        input_path = input_path.expanduser().resolve()
        output_dir = options.output_dir.expanduser().resolve()
        ensure_output_dir(output_dir)

        outputs: list[Path] = []
        for source in iter_supported_files(input_path, SUPPORTED_EXTENSIONS, options.recursive):
            output_path = output_dir / f"{source.stem}.md"
            markdown = self.convert_file(source, options)
            write_markdown(markdown, output_path, overwrite=options.overwrite)
            outputs.append(output_path)

        if not outputs:
            raise FileNotFoundError(f"No supported files found in: {input_path}")

        return outputs

    def convert_file(self, input_path: Path, options: ConvertOptions) -> str:
        converter = get_converter(input_path)
        return converter.convert(input_path, options)

# ============================================================
# file: src/markdown_converter/cli.py
# ============================================================
from pathlib import Path
import click
from markdown_converter.config import ConvertOptions
from markdown_converter.service import MarkdownConvertService


@click.command(context_settings={"help_option_names": ["-h", "--help"]})
@click.argument("input_path", type=click.Path(exists=True, path_type=Path))
@click.option("-o", "--output", "output_dir", type=click.Path(path_type=Path), default=Path("./output"), help="输出目录")
@click.option("--recursive", is_flag=True, help="递归转换目录中的文件")
@click.option("--assets", "export_assets", is_flag=True, help="导出图片等资源到 assets 目录，当前为预留接口")
@click.option("--no-overwrite", is_flag=True, help="不覆盖已有输出文件")
def main(input_path: Path, output_dir: Path, recursive: bool, export_assets: bool, no_overwrite: bool):
    """Convert documents and files to Markdown."""
    options = ConvertOptions(
        output_dir=output_dir,
        export_assets=export_assets,
        recursive=recursive,
        overwrite=not no_overwrite,
    )

    service = MarkdownConvertService()
    try:
        outputs = service.convert_path(input_path, options)
    except Exception as exc:
        raise click.ClickException(str(exc)) from exc

    for output in outputs:
        click.echo(f"Generated: {output}")


if __name__ == "__main__":
    main()

# ============================================================
# file: tests/test_utils.py
# ============================================================
from markdown_converter.utils import table_to_markdown


def test_table_to_markdown():
    rows = [
        ["Name", "Age"],
        ["Alice", 30],
        ["Bob", 28],
    ]
    result = table_to_markdown(rows)
    assert "| Name | Age |" in result
    assert "| Alice | 30 |" in result
