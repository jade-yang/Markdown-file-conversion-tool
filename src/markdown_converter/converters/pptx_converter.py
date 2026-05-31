from pathlib import Path
from pptx import Presentation
from markdown_converter.config import ConvertOptions
from markdown_converter.converters.base import BaseConverter
from markdown_converter.utils import normalize_text, table_to_markdown


class PptxConverter(BaseConverter):
    def convert(self, input_path: Path, options: ConvertOptions) -> str:
        prs = Presentation(str(input_path))
        parts: list[str] = []

        for slide_index, slide in enumerate(prs.slides, start=1):
            parts.append(f"# Slide {slide_index}")
            slide_text: list[str] = []

            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_text.append(shape.text.strip())

                if getattr(shape, "has_table", False):
                    rows = []
                    for row in shape.table.rows:
                        rows.append([cell.text for cell in row.cells])
                    table_md = table_to_markdown(rows)
                    if table_md:
                        slide_text.append(table_md)

            if slide_text:
                parts.append("\n\n".join(slide_text))

        return normalize_text("\n\n".join(parts))
